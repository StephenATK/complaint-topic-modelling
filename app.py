"""
app.py
Streamlit app for Group 8 - Customer Complaint Topic Modelling.

Run from PyCharm's TERMINAL (not the green Run button):
    streamlit run app.py
"""

import streamlit as st
import pandas as pd
import numpy as np
import joblib
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap
from wordcloud import WordCloud
import sys
import os
import io
import tempfile
from datetime import datetime
from fpdf import FPDF, XPos, YPos
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pyLDAvis

sys.path.append(os.path.join(os.path.dirname(__file__), "src"))
from preprocessing import preprocess_text

st.set_page_config(
    page_title="Customer Complaint Topic Explorer",
    page_icon="🗂️",
    layout="wide",
)

st.session_state.setdefault("topic_labels", {})
st.session_state.setdefault("app_theme", "Day")

# --- Fixed brand palette: used only by exported PDF/PPTX reports, which always
# render in the standard brand look regardless of the viewer's in-app theme choice. ---
NAVY = "#1b3865"
SKY = "#45c8f3"
ORANGE = "#f47a21"
LIGHT_GRAY = "#f4f6f9"
GRAY = "#5e5e5f"

TOPIC_COLORS = [
    NAVY, ORANGE, SKY,
    "#5c7ba3", "#f4a35c", "#8fdcf7",
    "#0f2540", "#c25e14", "#2f95b8",
    "#a4b8cd", GRAY,
]
TOPIC_COLORMAPS = [
    LinearSegmentedColormap.from_list("t", ["#FFFFFF", c]) for c in TOPIC_COLORS
]

# --- Two in-app themes, both built from the same 5 brand colors. Only affects the live UI. ---
THEMES = {
    "Day": {
        "app_bg": "#fdf6ec",
        "sidebar_bg": "#faeee0",
        "card_bg": "#faeee0",
        "card_border": "#f0ddc4",
        "text": "#1b3865",
        "text_secondary": "#8a6a45",
        "header_grad": "linear-gradient(135deg, #f47a21 0%, #e0601a 55%, #1b3865 130%)",
        "brand_icon_grad": "linear-gradient(135deg, #f47a21 0%, #1b3865 100%)",
        "nav_selected_bg": "#f47a21",
        "nav_selected_text": "#FFFFFF",
        "nav_hover_bg": "#f7e3cf",
        "wc_bg": "#fdf6ec",
        "grid_color": "#ecdcc4",
        "topic_colors": [
            "#f47a21", "#1b3865", "#45c8f3", "#e0601a", "#5c7ba3",
            "#ffb877", "#0f2540", "#8fdcf7", "#c25e14", "#2f95b8", "#a4b8cd",
        ],
    },
    "Night": {
        "app_bg": "#0d1826",
        "sidebar_bg": "#111f33",
        "card_bg": "#16283f",
        "card_border": "#26405f",
        "text": "#f2f5fa",
        "text_secondary": "#c3cedf",
        "header_grad": "linear-gradient(135deg, #0d1826 0%, #1b3865 55%, #45c8f3 130%)",
        "brand_icon_grad": "linear-gradient(135deg, #0d1826 0%, #45c8f3 100%)",
        "nav_selected_bg": "#45c8f3",
        "nav_selected_text": "#0d1826",
        "nav_hover_bg": "#1f3654",
        "wc_bg": "#16283f",
        "grid_color": "#2c4463",
        "topic_colors": [
            "#45c8f3", "#f47a21", "#8fdcf7", "#f4a35c", "#5c9fd6",
            "#ffb877", "#2f95b8", "#ffd1a3", "#a9c9e8", "#c25e14", "#dbe4f0",
        ],
    },
}

# --- Theme toggle: a single icon button showing the mode you'll switch TO.
# Rendered first so its value is known before the CSS below is built. ---
with st.sidebar:
    toggle_cols = st.columns([4, 1])
    with toggle_cols[1]:
        target_icon = "🌙" if st.session_state["app_theme"] == "Day" else "☀️"
        if st.button(target_icon, key="theme_toggle_btn", help="Switch day/night theme"):
            st.session_state["app_theme"] = (
                "Night" if st.session_state["app_theme"] == "Day" else "Day"
            )
            st.rerun()

selected_theme = st.session_state["app_theme"]
THEME = THEMES[selected_theme]
UI_TOPIC_COLORS = THEME["topic_colors"]
UI_TOPIC_COLORMAPS = [
    LinearSegmentedColormap.from_list("t", [THEME["wc_bg"], c]) for c in UI_TOPIC_COLORS
]

# --- Global CSS: typography, header banner, card-style metrics - all theme-driven ---
st.markdown(
    f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&display=swap');

    html, body, [class*="css"] {{
        font-family: 'Inter', sans-serif;
    }}

    .stApp {{
        background-color: {THEME["app_bg"]};
    }}
    .stApp, .stApp p, .stApp li, .stApp span, .stApp label {{
        color: {THEME["text"]};
    }}

    /* Streamlit's own widgets (labels, captions, expanders, file uploader) use
       data-testid attributes rather than semantic tags, so the broad rule above
       doesn't reach them - covered explicitly here. */
    .stApp [data-testid="stWidgetLabel"] p,
    .stApp [data-testid="stMarkdownContainer"] p,
    .stApp [data-testid="stMarkdownContainer"] li,
    .stApp [data-testid="stMarkdownContainer"] strong {{
        color: {THEME["text"]} !important;
    }}
    .stApp [data-testid="stCaptionContainer"],
    .stApp [data-testid="stCaptionContainer"] p {{
        color: {THEME["text_secondary"]} !important;
    }}
    .stApp [data-testid="stExpander"] {{
        background-color: {THEME["card_bg"]};
        border: 1px solid {THEME["card_border"]};
        border-radius: 10px;
    }}
    .stApp [data-testid="stExpander"] summary p,
    .stApp [data-testid="stExpander"] summary span {{
        color: {THEME["text"]} !important;
    }}
    .stApp [data-testid="stFileUploaderDropzone"] {{
        background-color: {THEME["card_bg"]};
        border: 1px solid {THEME["card_border"]};
    }}
    .stApp [data-testid="stFileUploaderDropzone"] div,
    .stApp [data-testid="stFileUploaderDropzone"] span,
    .stApp [data-testid="stFileUploaderDropzoneInstructions"] div,
    .stApp [data-testid="stFileUploaderDropzoneInstructions"] span {{
        color: {THEME["text_secondary"]} !important;
    }}
    .stApp [data-testid="stFileUploaderFile"] span,
    .stApp [data-testid="stFileUploaderFileName"] {{
        color: {THEME["text"]} !important;
    }}
    .stApp [data-testid="stDataFrame"] {{
        border: 1px solid {THEME["card_border"]};
        border-radius: 8px;
    }}
    .stApp [data-testid="stTable"] table,
    .stApp [data-testid="stTable"] th,
    .stApp [data-testid="stTable"] td {{
        color: {THEME["text"]} !important;
        background-color: {THEME["card_bg"]} !important;
        border-color: {THEME["card_border"]} !important;
    }}
    .stApp [data-testid="stAlertContentInfo"],
    .stApp [data-testid="stAlertContentInfo"] p,
    .stApp [data-testid="stAlertContentSuccess"],
    .stApp [data-testid="stAlertContentSuccess"] p,
    .stApp [data-testid="stAlertContentWarning"],
    .stApp [data-testid="stAlertContentWarning"] p,
    .stApp [data-testid="stAlertContentError"],
    .stApp [data-testid="stAlertContentError"] p {{
        color: #1b3865 !important;
    }}
    .stApp hr {{
        border-color: {THEME["card_border"]} !important;
    }}
    .stApp [data-testid="stTextInput"] input,
    .stApp [data-testid="stSelectbox"] div[data-baseweb="select"] {{
        color: #1b3865 !important;
    }}

    .app-header {{
        padding: 1.75rem 2rem;
        border-radius: 12px;
        background: {THEME["header_grad"]};
        color: white;
        margin-bottom: 1.5rem;
    }}
    .app-header h1 {{
        margin: 0;
        font-size: 1.9rem;
        font-weight: 700;
        color: white;
    }}
    .app-header p {{
        margin: 0.25rem 0 0 0;
        font-size: 0.95rem;
        opacity: 0.92;
        color: white;
    }}
    /* Wins over the stMarkdownContainer text-color fix below, since the header
       is rendered via st.markdown too and would otherwise inherit the theme's
       body text color instead of staying white-on-gradient. */
    .stApp [data-testid="stMarkdownContainer"] .app-header h1,
    .stApp [data-testid="stMarkdownContainer"] .app-header p {{
        color: #FFFFFF !important;
        opacity: 0.92;
    }}
    .stApp [data-testid="stMarkdownContainer"] .app-header h1 {{
        opacity: 1;
    }}

    div[data-testid="stMetric"] {{
        background-color: {THEME["card_bg"]};
        border: 1px solid {THEME["card_border"]};
        border-left: 3px solid {THEME["nav_selected_bg"]};
        border-radius: 10px;
        padding: 0.9rem 1rem;
    }}
    div[data-testid="stMetric"] label {{
        color: {THEME["text_secondary"]} !important;
    }}
    div[data-testid="stMetricValue"] {{
        color: {THEME["text"]} !important;
    }}

    h2, h3, h4 {{
        color: {THEME["text"]};
        font-weight: 700;
    }}

    section[data-testid="stSidebar"] {{
        background-color: {THEME["sidebar_bg"]};
        border-right: 1px solid {THEME["card_border"]};
    }}

    .sidebar-brand {{
        display: flex;
        align-items: center;
        gap: 0.65rem;
        margin: 0.25rem 0 1.5rem 0;
    }}
    .sidebar-brand-icon {{
        width: 42px;
        height: 42px;
        border-radius: 11px;
        background: {THEME["brand_icon_grad"]};
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 1.35rem;
        flex-shrink: 0;
        box-shadow: 0 2px 6px rgba(0, 0, 0, 0.25);
    }}
    .sidebar-brand-text {{
        font-weight: 700;
        font-size: 1rem;
        line-height: 1.25;
        color: {THEME["text"]};
    }}
    .sidebar-brand-text small {{
        display: block;
        font-weight: 400;
        font-size: 0.72rem;
        color: {THEME["text_secondary"]};
        letter-spacing: 0.03em;
        text-transform: uppercase;
    }}

    section[data-testid="stSidebar"] div[role="radiogroup"] {{
        gap: 0.25rem;
    }}
    section[data-testid="stSidebar"] div[role="radiogroup"] label {{
        padding: 0.6rem 0.85rem;
        border-radius: 8px;
        width: 100%;
        transition: background-color 0.15s ease;
    }}
    section[data-testid="stSidebar"] div[role="radiogroup"] label p {{
        font-size: 1.02rem;
        color: {THEME["text"]};
    }}
    section[data-testid="stSidebar"] div[role="radiogroup"] label:hover {{
        background-color: {THEME["nav_hover_bg"]};
    }}
    section[data-testid="stSidebar"] div[role="radiogroup"] label:has(input:checked) {{
        background-color: {THEME["nav_selected_bg"]};
    }}
    section[data-testid="stSidebar"] div[role="radiogroup"] label:has(input:checked) p {{
        color: {THEME["nav_selected_text"]} !important;
        font-weight: 600;
    }}
    section[data-testid="stSidebar"] div[role="radiogroup"] svg {{
        display: none;
    }}

    .pipeline-step {{
        background-color: {THEME["card_bg"]} !important;
        border-left: 4px solid {THEME["nav_selected_bg"]} !important;
    }}
    .pipeline-step-title {{
        color: {THEME["text"]} !important;
    }}
    .pipeline-step-desc {{
        color: {THEME["text_secondary"]} !important;
    }}
    </style>
    """,
    unsafe_allow_html=True,
)


def page_header(title, subtitle):
    """Consistent gradient banner header used at the top of every page.
    Colors are set as inline styles with !important - nothing in the external
    stylesheet (including the broad theme text-color rules) can override an
    inline !important, so this is guaranteed correct regardless of Streamlit's
    internal HTML structure, which isn't something we can rely on assumptions about."""
    st.markdown(
        f"""
        <div class="app-header">
            <h1 style="color:#FFFFFF !important; margin:0; font-size:1.9rem; font-weight:700;">{title}</h1>
            <p style="color:#FFFFFF !important; opacity:0.92; margin:0.25rem 0 0 0; font-size:0.95rem;">{subtitle}</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_wordcloud(freq_or_text, colormap="Blues", bg_color="white"):
    """Accepts either a raw text string or a {word: weight} dict and renders a word cloud.
    bg_color should match the current theme so there's no white box on dark themes."""
    wc = WordCloud(
        width=1000, height=400, background_color=bg_color,
        colormap=colormap, max_words=100, prefer_horizontal=0.95,
    )
    if isinstance(freq_or_text, dict):
        wc = wc.generate_from_frequencies(freq_or_text)
    else:
        wc = wc.generate(freq_or_text)
    fig, ax = plt.subplots(figsize=(10, 4))
    fig.patch.set_facecolor(bg_color)
    ax.imshow(wc, interpolation="bilinear")
    ax.axis("off")
    fig.tight_layout(pad=0)
    return fig


def themed_axes(fig, ax):
    """Applies the current theme's colors to a matplotlib fig/ax pair -
    background, tick labels, axis labels, title, and spines - so every
    custom chart matches whichever theme is active."""
    fig.patch.set_facecolor(THEME["card_bg"])
    ax.set_facecolor(THEME["card_bg"])
    ax.tick_params(colors=THEME["text_secondary"])
    ax.xaxis.label.set_color(THEME["text"])
    ax.yaxis.label.set_color(THEME["text"])
    ax.title.set_color(THEME["text"])
    for spine in ax.spines.values():
        spine.set_color(THEME["grid_color"])
    return fig, ax


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


def pdf_safe(text):
    """fpdf2's built-in Helvetica font only supports Latin-1, but the app's text -
    including user-typed topic names - can contain em-dashes, curly quotes, or
    emoji. Map common punctuation to ASCII equivalents and fall back to '?' for
    anything else unsupported, rather than letting the PDF export crash."""
    replacements = {
        "\u2014": "-", "\u2013": "-", "\u2018": "'", "\u2019": "'",
        "\u201c": '"', "\u201d": '"', "\u2026": "...",
    }
    for uni_char, ascii_char in replacements.items():
        text = text.replace(uni_char, ascii_char)
    return text.encode("latin-1", "replace").decode("latin-1")


def get_topic_label(idx):
    """Returns the human-assigned name for a topic if one's been set (via the
    Topic Explorer page's 'Name Your Topics' panel), otherwise falls back to 'Topic N'."""
    return st.session_state["topic_labels"].get(idx, f"Topic {idx}")


def generate_executive_summary(dist_series, total_docs, model_choice):
    """Auto-written narrative paragraph summarizing the topic distribution -
    reusable in-app, in the PDF report, and in the slide deck."""
    top_idx = int(dist_series.idxmax())
    top_count = int(dist_series.max())
    top_pct = (top_count / total_docs * 100) if total_docs else 0
    n_topics = len(dist_series)
    label = get_topic_label(top_idx)
    return (
        f"Across {total_docs:,} complaints analyzed with {model_choice}, the largest topic — "
        f"\"{label}\" — accounts for {top_pct:.0f}% of documents ({top_count:,} complaints). "
        f"The remaining {total_docs - top_count:,} complaints are distributed across the other "
        f"{n_topics - 1} topics, suggesting a mix of dominant and long-tail complaint themes."
    )


# A small built-in sample so the "Analyze & Report" page can be demoed instantly,
# without needing a live file upload during a presentation.
SAMPLE_COMPLAINTS = [
    "A debt collector has been calling me multiple times a day, even after I asked them to stop contacting me at work.",
    "I received a collection notice for a debt that is not mine, and I have already disputed this with the company twice.",
    "The collection agency threatened to sue me and garnish my wages if I did not pay immediately.",
    "I paid this debt in full last year but the collector is still reporting it as unpaid on my credit report.",
    "My credit report shows an account that I never opened, and I believe I am a victim of identity theft.",
    "I disputed an error on my credit report three months ago and it still has not been corrected.",
    "The company keeps calling my family members and neighbors about my debt, which feels like harassment.",
    "I asked for written validation of this debt and never received anything, yet they continue to call.",
    "My credit score dropped significantly after a collection account appeared that I do not recognize.",
    "The collector used abusive language and made threats during a phone call about a medical bill.",
    "I was never notified about this debt before it appeared on my credit report as a collection account.",
    "I have been trying to reach the company for weeks to correct an error in the amount owed.",
    "This account was already discharged in bankruptcy but is still being pursued by a collection agency.",
    "The debt collector contacted my employer directly, which I believe is against the rules.",
    "I keep receiving calls early in the morning and late at night about a debt I already settled.",
]


def build_pdf_report(new_df, assignments, topic_dist, model, feature_names, model_choice, source_name):
    """Builds a branded PDF summarizing topic assignments for an uploaded file.
    Returns raw PDF bytes, ready for st.download_button."""
    n_topics = model.components_.shape[0]
    dist = pd.Series(assignments).value_counts().reindex(range(n_topics), fill_value=0)
    summary_text = generate_executive_summary(dist, len(new_df), model_choice)

    # --- Chart: documents per topic, in brand colors, labeled with custom topic names ---
    fig, ax = plt.subplots(figsize=(7, 3.2))
    bar_colors = [TOPIC_COLORS[i % len(TOPIC_COLORS)] for i in dist.index]
    labels_x = [get_topic_label(i) for i in dist.index]
    ax.bar(labels_x, dist.values, color=bar_colors)
    ax.set_ylabel("Documents")
    ax.set_title("Documents per Topic")
    ax.spines[["top", "right"]].set_visible(False)
    plt.xticks(rotation=25, ha="right", fontsize=8)
    fig.tight_layout()

    with tempfile.TemporaryDirectory() as tmpdir:
        chart_path = os.path.join(tmpdir, "chart.png")
        fig.savefig(chart_path, dpi=150)
        plt.close(fig)

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        # --- Header ---
        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(*hex_to_rgb(NAVY))
        pdf.cell(0, 12, "Customer Complaint Topic Report",
                 new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*hex_to_rgb(GRAY))
        pdf.cell(0, 8,
                 pdf_safe(f"Source file: {source_name}  |  Model: {model_choice}  |  "
                          f"Generated {datetime.now().strftime('%B %d, %Y %H:%M')}"),
                 new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

        pdf.set_draw_color(*hex_to_rgb(SKY))
        pdf.set_line_width(0.8)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(6)

        # --- Summary ---
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(*hex_to_rgb(NAVY))
        pdf.cell(0, 8, "Executive Summary", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(40, 40, 40)
        pdf.multi_cell(0, 6, pdf_safe(summary_text))
        pdf.ln(2)

        pdf.image(chart_path, x=15, w=180)
        pdf.ln(4)

        # --- Per-topic keyword breakdown ---
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(*hex_to_rgb(NAVY))
        pdf.cell(0, 8, "Topic Keywords", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        for i in range(n_topics):
            weights = model.components_[i]
            top_idx = weights.argsort()[:-11:-1]
            words = ", ".join(feature_names[j] for j in top_idx)
            doc_count = int(dist.get(i, 0))

            pdf.set_font("Helvetica", "B", 10)
            pdf.set_text_color(*hex_to_rgb(TOPIC_COLORS[i % len(TOPIC_COLORS)]))
            pdf.multi_cell(0, 5.5, pdf_safe(f"{get_topic_label(i)}  ({doc_count} documents)"),
                            new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("Helvetica", "", 9)
            pdf.set_text_color(60, 60, 60)
            pdf.multi_cell(0, 5.5, pdf_safe(words), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(1.5)

        return bytes(pdf.output())


def _pptx_slide_header(slide, title, color_rgb, prs):
    """Colored header bar + title text, reused across every content slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color_rgb
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(11.5), Inches(0.6))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(26)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def build_pptx_report(new_df, assignments, topic_dist, model, feature_names, model_choice, source_name):
    """Builds a branded slide deck (title, distribution chart, executive summary,
    one slide per topic) from the same analysis results as the PDF report."""
    n_topics = model.components_.shape[0]
    dist = pd.Series(assignments).value_counts().reindex(range(n_topics), fill_value=0)
    summary_text = generate_executive_summary(dist, len(new_df), model_choice)

    fig, ax = plt.subplots(figsize=(8, 4.2))
    bar_colors = [TOPIC_COLORS[i % len(TOPIC_COLORS)] for i in dist.index]
    labels_x = [get_topic_label(i) for i in dist.index]
    ax.bar(labels_x, dist.values, color=bar_colors)
    ax.set_ylabel("Documents")
    ax.set_title("Documents per Topic")
    ax.spines[["top", "right"]].set_visible(False)
    plt.xticks(rotation=25, ha="right")
    fig.tight_layout()

    with tempfile.TemporaryDirectory() as tmpdir:
        chart_path = os.path.join(tmpdir, "chart.png")
        fig.savefig(chart_path, dpi=150)
        plt.close(fig)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        navy_rgb = RGBColor(*hex_to_rgb(NAVY))
        gray_rgb = RGBColor(*hex_to_rgb(GRAY))

        # --- Title slide ---
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = navy_rgb
        bg.line.fill.background()

        title_tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.5))
        title_tb.text_frame.text = "Customer Complaint Topic Report"
        title_tb.text_frame.paragraphs[0].font.size = Pt(40)
        title_tb.text_frame.paragraphs[0].font.bold = True
        title_tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        sub_tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8))
        sub_tb.text_frame.text = (
            f"Source: {source_name}  ·  Model: {model_choice}  ·  "
            f"{datetime.now().strftime('%B %d, %Y')}"
        )
        sub_tb.text_frame.paragraphs[0].font.size = Pt(16)
        sub_tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb(SKY))

        # --- Executive summary slide ---
        slide_sum = prs.slides.add_slide(blank_layout)
        _pptx_slide_header(slide_sum, "Executive Summary", navy_rgb, prs)
        sum_tb = slide_sum.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5))
        sum_tb.text_frame.word_wrap = True
        sum_tb.text_frame.text = summary_text
        sum_tb.text_frame.paragraphs[0].font.size = Pt(20)
        sum_tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(40, 40, 40)

        # --- Distribution chart slide ---
        slide_dist = prs.slides.add_slide(blank_layout)
        _pptx_slide_header(slide_dist, "Topic Distribution", navy_rgb, prs)
        slide_dist.shapes.add_picture(chart_path, Inches(0.9), Inches(1.3), width=Inches(11.5))

        # --- One slide per topic ---
        for i in range(n_topics):
            weights = model.components_[i]
            top_idx = weights.argsort()[:-11:-1]
            words = ", ".join(feature_names[j] for j in top_idx)
            doc_count = int(dist.get(i, 0))

            slide_k = prs.slides.add_slide(blank_layout)
            color = RGBColor(*hex_to_rgb(TOPIC_COLORS[i % len(TOPIC_COLORS)]))
            _pptx_slide_header(slide_k, get_topic_label(i), color, prs)

            count_tb = slide_k.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.6))
            count_tb.text_frame.text = f"{doc_count} documents"
            count_tb.text_frame.paragraphs[0].font.size = Pt(18)
            count_tb.text_frame.paragraphs[0].font.color.rgb = gray_rgb

            words_tb = slide_k.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(3.8))
            words_tb.text_frame.word_wrap = True
            words_tb.text_frame.text = words
            words_tb.text_frame.paragraphs[0].font.size = Pt(24)
            words_tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(40, 40, 40)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


@st.cache_data
def load_data():
    df = pd.read_csv("data/complaints_clean.csv")
    return df


@st.cache_resource
def load_models():
    tfidf_vectorizer = joblib.load("data/tfidf_vectorizer.joblib")
    count_vectorizer = joblib.load("data/count_vectorizer.joblib")
    lda_model = joblib.load("data/lda_model.joblib")
    nmf_model = joblib.load("data/nmf_model.joblib")
    lda_doc_topic = np.load("data/lda_doc_topic.npy")
    nmf_doc_topic = np.load("data/nmf_doc_topic.npy")
    return tfidf_vectorizer, count_vectorizer, lda_model, nmf_model, lda_doc_topic, nmf_doc_topic


def get_top_words(model, feature_names, n_top=10):
    topics = []
    for topic_weights in model.components_:
        top_indices = topic_weights.argsort()[: -n_top - 1 : -1]
        topics.append([feature_names[i] for i in top_indices])
    return topics


@st.cache_data(show_spinner=False)
def build_pyldavis_html(_model, feature_names, doc_topic, _doc_term_matrix, model_choice):
    """Builds an interactive pyLDAvis topic map. Works for both LDA and NMF -
    NMF's outputs aren't proper probability distributions by default, so they're
    normalized to sum to 1 before handing them to pyLDAvis (which expects
    genuine distributions). Leading-underscore args are excluded from Streamlit's
    cache key since models/sparse matrices aren't reliably hashable; feature_names,
    doc_topic, and model_choice are hashable and determine cache validity."""
    topic_term = _model.components_
    topic_term_dists = topic_term / topic_term.sum(axis=1, keepdims=True)
    doc_topic_dists = doc_topic / doc_topic.sum(axis=1, keepdims=True)
    doc_lengths = np.asarray(_doc_term_matrix.sum(axis=1)).flatten()
    term_frequency = np.asarray(_doc_term_matrix.sum(axis=0)).flatten()

    vis_data = pyLDAvis.prepare(
        topic_term_dists=topic_term_dists,
        doc_topic_dists=doc_topic_dists,
        doc_lengths=doc_lengths,
        vocab=feature_names,
        term_frequency=term_frequency,
        sort_topics=False,
    )
    return pyLDAvis.prepared_data_to_html(vis_data)


# --- Sidebar navigation ---
st.sidebar.markdown(
    """
    <div class="sidebar-brand">
        <div class="sidebar-brand-icon">📋</div>
        <div class="sidebar-brand-text">
            Complaint Topic<br>Modelling
            <small>Group 8 · MSBA610</small>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

PAGE_NAMES = ["How It Works", "Corpus Overview", "Preprocessing Demo", "Topic Explorer", "Model Comparison", "Try It Yourself", "Analyze & Report"]
PAGE_ICONS = {
    "How It Works": "🧭",
    "Corpus Overview": "📊",
    "Preprocessing Demo": "🧹",
    "Topic Explorer": "🔍",
    "Model Comparison": "⚖️",
    "Try It Yourself": "✍️",
    "Analyze & Report": "📤",
}
nav_choice = st.sidebar.radio(
    "Navigate",
    [f"{PAGE_ICONS[name]}  {name}" for name in PAGE_NAMES],
    label_visibility="collapsed",
)
page = nav_choice.split("  ", 1)[1]

data_ready = os.path.exists("data/complaints_clean.csv")
models_ready = os.path.exists("data/lda_model.joblib")

# ============================================================
# PAGE: How It Works (viewable even before the pipeline has been run)
# ============================================================
if page == "How It Works":
    page_header("How It Works", "The pipeline behind this app, in five steps.")

    st.markdown(
        """
        <style>
        .pipeline-step {
            display: flex;
            gap: 1rem;
            padding: 1rem 1.2rem;
            border-radius: 10px;
            background-color: #f4f6f9;
            border-left: 4px solid #45c8f3;
            margin-bottom: 0.85rem;
        }
        .pipeline-step-icon {
            font-size: 1.6rem;
            flex-shrink: 0;
        }
        .pipeline-step-title {
            font-weight: 700;
            color: #1b3865;
            margin-bottom: 0.15rem;
        }
        .pipeline-step-desc {
            color: #3a3a3a;
            font-size: 0.92rem;
            line-height: 1.4;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    steps = [
        ("🧹", "Clean & Preprocess",
         "Raw complaint text is stripped of redacted placeholders, dollar amounts, URLs, and "
         "punctuation, then tokenized, stripped of stopwords, and lemmatized so 'calling', "
         "'called', and 'calls' all collapse to one meaningful term."),
        ("🔢", "Represent as Numbers",
         "Cleaned text is converted into two numeric forms: TF-IDF weighted vectors for NMF, "
         "and raw word counts for LDA — matching how each algorithm is designed to work."),
        ("🧩", "Discover Topics",
         "Two unsupervised models — Latent Dirichlet Allocation (LDA) and Non-negative Matrix "
         "Factorization (NMF) — each find a set of latent topics purely from word co-occurrence "
         "patterns, with no labels required."),
        ("📏", "Evaluate & Tune",
         "Topic coherence (do a topic's top words actually relate to each other?) and topic "
         "diversity (are topics distinct, not redundant?) are measured across a range of topic "
         "counts to pick the best-performing model and number of topics."),
        ("🗂️", "Explore & Deploy",
         "The winning model powers this app — browse topics, test new text, or upload a file "
         "for live analysis and a downloadable report, all running the model trained during "
         "development."),
    ]
    for icon, title, desc in steps:
        st.markdown(
            f"""
            <div class="pipeline-step">
                <div class="pipeline-step-icon">{icon}</div>
                <div>
                    <div class="pipeline-step-title">{title}</div>
                    <div class="pipeline-step-desc">{desc}</div>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    st.info(
        "Use the sidebar to explore each stage: **Corpus Overview** and **Preprocessing Demo** "
        "cover steps 1-2, **Topic Explorer** and **Model Comparison** cover steps 3-4, and "
        "**Try It Yourself** / **Analyze & Report** cover step 5."
    )
    st.stop()

if not data_ready:
    st.warning(
        "No processed data found yet. Run `python src/preprocessing.py` from the terminal first, "
        "then `python src/modeling.py` to train the topic models."
    )
    st.stop()

df = load_data()

# ============================================================
# PAGE: Corpus Overview
# ============================================================
if page == "Corpus Overview":
    page_header("Corpus Overview", "A first look at the complaint dataset before any modelling.")

    col1, col2, col3 = st.columns(3)
    col1.metric("Total complaints", f"{len(df):,}")
    col2.metric("Avg. words / complaint", f"{df['clean_text'].str.split().str.len().mean():.0f}")
    col3.metric("Vocabulary size (approx.)", f"{len(set(' '.join(df['clean_text']).split())):,}")

    st.subheader("Complaint length distribution")
    fig, ax = plt.subplots(figsize=(8, 3))
    themed_axes(fig, ax)
    df["clean_text"].str.split().str.len().hist(bins=40, ax=ax, color=UI_TOPIC_COLORS[0])
    ax.set_xlabel("Word count")
    ax.set_ylabel("Number of complaints")
    ax.spines[["top", "right"]].set_visible(False)
    st.pyplot(fig)

    col_a, col_b = st.columns(2)
    with col_a:
        st.subheader("Top words in the cleaned corpus")
        all_words = " ".join(df["clean_text"]).split()
        top_words = pd.Series(all_words).value_counts().head(15)
        fig_w, ax_w = plt.subplots(figsize=(6, 4.2))
        themed_axes(fig_w, ax_w)
        ax_w.barh(top_words.index[::-1], top_words.values[::-1], color=UI_TOPIC_COLORS[1])
        ax_w.spines[["top", "right"]].set_visible(False)
        st.pyplot(fig_w)

    with col_b:
        st.subheader("☁️ Word Cloud")
        st.pyplot(render_wordcloud(
            " ".join(df["clean_text"]), colormap=UI_TOPIC_COLORMAPS[2], bg_color=THEME["wc_bg"],
        ))

# ============================================================
# PAGE: Preprocessing Demo
# ============================================================
elif page == "Preprocessing Demo":
    page_header("Preprocessing Demo", "See the cleaning pipeline applied to raw text, live.")

    sample = st.text_area(
        "Complaint text",
        value="I called XXXX regarding my $500.00 charge and they never responded to my email at test@example.com!!",
        height=120,
    )
    if sample:
        st.subheader("Cleaned output")
        st.code(preprocess_text(sample), language=None)

# ============================================================
# PAGE: Topic Explorer
# ============================================================
elif page == "Topic Explorer":
    page_header("Topic Explorer", "Browse each discovered topic and its defining vocabulary.")

    if not models_ready:
        st.warning("Run `python src/modeling.py` from the terminal first to train the models.")
        st.stop()

    tfidf_vectorizer, count_vectorizer, lda_model, nmf_model, lda_doc_topic, nmf_doc_topic = load_models()

    model_choice = st.selectbox("Choose a model", ["LDA", "NMF"])

    if model_choice == "LDA":
        model, feature_names = lda_model, count_vectorizer.get_feature_names_out()
        doc_topic = lda_doc_topic
    else:
        model, feature_names = nmf_model, tfidf_vectorizer.get_feature_names_out()
        doc_topic = nmf_doc_topic

    n_topics = model.components_.shape[0]

    with st.expander("✏️ Name Your Topics", expanded=False):
        st.caption(
            "Replace generic topic numbers with meaningful names once you've reviewed the "
            "keywords below — these names carry through every chart, the PDF report, and the slide deck."
        )
        label_cols = st.columns(2)
        for i in range(n_topics):
            col = label_cols[i % 2]
            with col:
                new_label = st.text_input(
                    f"Topic {i}", value=get_topic_label(i), key=f"label_input_{i}",
                )
                st.session_state["topic_labels"][i] = new_label

    st.subheader(f"{model_choice} — Topic Word Clouds")
    st.caption("Word size reflects how strongly that word defines the topic.")

    # Two topics per row, each with its own color-coded word cloud
    for row_start in range(0, n_topics, 2):
        cols = st.columns(2)
        for offset, col in enumerate(cols):
            topic_idx = row_start + offset
            if topic_idx >= n_topics:
                continue
            topic_weights = model.components_[topic_idx]
            top_indices = topic_weights.argsort()[:-31:-1]
            freq = {feature_names[i]: float(topic_weights[i]) for i in top_indices}
            color = UI_TOPIC_COLORS[topic_idx % len(UI_TOPIC_COLORS)]
            colormap = UI_TOPIC_COLORMAPS[topic_idx % len(UI_TOPIC_COLORMAPS)]
            with col:
                st.markdown(
                    f"<h4 style='color:{color}; margin-bottom:0.2rem;'>{get_topic_label(topic_idx)}</h4>",
                    unsafe_allow_html=True,
                )
                st.pyplot(render_wordcloud(freq, colormap=colormap, bg_color=THEME["wc_bg"]))
                top5 = sorted(freq, key=freq.get, reverse=True)[:5]
                st.caption(", ".join(top5))

    st.subheader("Documents per topic")
    assignments = np.argmax(doc_topic, axis=1)
    dist = pd.Series(assignments).value_counts().sort_index()
    fig_d, ax_d = plt.subplots(figsize=(8, 3.2))
    themed_axes(fig_d, ax_d)
    bar_colors = [UI_TOPIC_COLORS[i % len(UI_TOPIC_COLORS)] for i in dist.index]
    ax_d.bar([get_topic_label(i) for i in dist.index], dist.values, color=bar_colors)
    ax_d.set_ylabel("Documents")
    plt.setp(ax_d.get_xticklabels(), rotation=25, ha="right")
    ax_d.spines[["top", "right"]].set_visible(False)
    st.pyplot(fig_d)

    st.info(generate_executive_summary(dist, len(df), model_choice))

    st.subheader("🗺️ Interactive Topic Map")
    st.caption(
        "Bubble size shows how prevalent a topic is; bubble distance shows how related topics "
        "are to each other. Click a bubble to see its top terms, or hover over a bar to see "
        "which topics use that term most."
    )
    with st.spinner("Building interactive map (this can take a few seconds)..."):
        source_matrix = count_vectorizer.transform(df["clean_text"]) if model_choice == "LDA" \
            else tfidf_vectorizer.transform(df["clean_text"])
        vis_html = build_pyldavis_html(model, feature_names, doc_topic, source_matrix, model_choice)
    st.iframe(vis_html, height=800)

# ============================================================
# PAGE: Model Comparison
# ============================================================
elif page == "Model Comparison":
    page_header("Model Comparison", "LDA vs NMF, side by side on coherence and diversity.")

    sweep_path = "data/n_topics_sweep.csv"
    if os.path.exists(sweep_path):
        sweep_df = pd.read_csv(sweep_path)
        fig_c, ax_c = plt.subplots(figsize=(8, 3.5))
        themed_axes(fig_c, ax_c)
        ax_c.plot(sweep_df["n_topics"], sweep_df["lda_coherence"], marker="o",
                  color=UI_TOPIC_COLORS[1], label="LDA", linewidth=2)
        ax_c.plot(sweep_df["n_topics"], sweep_df["nmf_coherence"], marker="o",
                  color=UI_TOPIC_COLORS[0], label="NMF", linewidth=2)
        ax_c.set_xlabel("Number of topics")
        ax_c.set_ylabel("Coherence (c_v)")
        legend = ax_c.legend(facecolor=THEME["card_bg"], edgecolor=THEME["grid_color"])
        for text in legend.get_texts():
            text.set_color(THEME["text"])
        ax_c.spines[["top", "right"]].set_visible(False)
        st.pyplot(fig_c)

        best_lda = sweep_df.loc[sweep_df["lda_coherence"].idxmax()]
        best_nmf = sweep_df.loc[sweep_df["nmf_coherence"].idxmax()]

        col1, col2 = st.columns(2)
        col1.metric("Best LDA coherence", f"{best_lda['lda_coherence']:.3f}", f"at n_topics={int(best_lda['n_topics'])}")
        col2.metric("Best NMF coherence", f"{best_nmf['nmf_coherence']:.3f}", f"at n_topics={int(best_nmf['n_topics'])}")
    else:
        st.info("Run `python src/evaluation.py` first to generate data/n_topics_sweep.csv.")

    st.subheader("Summary")
    comparison_df = pd.DataFrame({
        "Metric": ["Topic Coherence (c_v)", "Topic Diversity"],
        "LDA": ["-", "-"],
        "NMF": ["-", "-"],
    })
    st.table(comparison_df)
    st.caption(
        "Fill the Topic Diversity row in with values from your own diversity calculation, "
        "and add 1-2 sentences below on which model won and why."
    )

# ============================================================
# PAGE: Try It Yourself
# ============================================================
elif page == "Try It Yourself":
    page_header("Try It Yourself", "Paste a new complaint and see which topic it's assigned to.")

    if not models_ready:
        st.warning("Run `python src/modeling.py` from the terminal first to train the models.")
        st.stop()

    tfidf_vectorizer, count_vectorizer, lda_model, nmf_model, lda_doc_topic, nmf_doc_topic = load_models()
    model_choice = st.selectbox("Model", ["LDA", "NMF"])

    user_text = st.text_area("Paste a complaint here", height=150)
    if st.button("Classify Topic") and user_text.strip():
        cleaned = preprocess_text(user_text)
        if model_choice == "LDA":
            vec = count_vectorizer.transform([cleaned])
            topic_dist = lda_model.transform(vec)[0]
            topics = get_top_words(lda_model, count_vectorizer.get_feature_names_out())
        else:
            vec = tfidf_vectorizer.transform([cleaned])
            topic_dist = nmf_model.transform(vec)[0]
            topics = get_top_words(nmf_model, tfidf_vectorizer.get_feature_names_out())

        best_topic = int(np.argmax(topic_dist))
        st.success(f"Assigned to **{get_topic_label(best_topic)}** (confidence: {topic_dist[best_topic]:.2f})")
        st.write(f"Top words: {', '.join(topics[best_topic])}")

# ============================================================
# PAGE: Analyze & Report
# ============================================================
elif page == "Analyze & Report":
    page_header("Analyze & Report", "Upload a new file (or try a sample), run it through the topic models, and download a report.")

    if not models_ready:
        st.warning("Run `python src/modeling.py` from the terminal first to train the models.")
        st.stop()

    def run_upload_analysis(raw_texts, base_df, model_choice, source_name):
        tfidf_vectorizer, count_vectorizer, lda_model, nmf_model, _, _ = load_models()
        cleaned = [preprocess_text(t) for t in raw_texts]

        if model_choice == "NMF":
            vecs = tfidf_vectorizer.transform(cleaned)
            topic_dist = nmf_model.transform(vecs)
        else:
            vecs = count_vectorizer.transform(cleaned)
            topic_dist = lda_model.transform(vecs)

        assignments = np.argmax(topic_dist, axis=1)
        result_df = base_df.copy()
        result_df["assigned_topic"] = [get_topic_label(i) for i in assignments]
        result_df["confidence"] = topic_dist[np.arange(len(assignments)), assignments].round(3)

        st.session_state["upload_analysis"] = {
            "df": result_df,
            "assignments": assignments,
            "topic_dist": topic_dist,
            "model_choice": model_choice,
            "source_name": source_name,
        }
        st.session_state.pop("report_bytes", None)
        st.session_state.pop("deck_bytes", None)

    col_upload, col_sample = st.columns([3, 1])
    with col_upload:
        uploaded_file = st.file_uploader(
            "Upload a CSV or plain-text file",
            type=["csv", "txt"],
            help="CSV: pick which column holds the complaint text. TXT: one complaint per line.",
        )
    with col_sample:
        st.markdown("<div style='height: 1.8rem'></div>", unsafe_allow_html=True)
        sample_clicked = st.button("🔁 Try a Sample", width="stretch")

    model_choice = st.selectbox("Model to use for topic assignment", ["NMF", "LDA"], key="upload_model_choice")

    if sample_clicked:
        with st.spinner("Running the sample through the topic model..."):
            sample_df = pd.DataFrame({"text": SAMPLE_COMPLAINTS})
            run_upload_analysis(SAMPLE_COMPLAINTS, sample_df, model_choice, "sample_complaints.csv")

    elif uploaded_file is not None:
        if uploaded_file.name.endswith(".csv"):
            upload_df = pd.read_csv(uploaded_file)
            text_col = st.selectbox("Which column contains the complaint text?", upload_df.columns)
            raw_texts = upload_df[text_col].astype(str).tolist()
        else:
            raw_texts = [
                line.strip() for line in uploaded_file.read().decode("utf-8").splitlines()
                if line.strip()
            ]
            upload_df = pd.DataFrame({"text": raw_texts})

        st.success(f"Loaded {len(raw_texts)} documents from **{uploaded_file.name}**.")

        if st.button("Run Topic Analysis", type="primary"):
            with st.spinner("Cleaning, preprocessing, and assigning topics..."):
                run_upload_analysis(raw_texts, upload_df, model_choice, uploaded_file.name)

    if "upload_analysis" in st.session_state:
        results = st.session_state["upload_analysis"]
        result_df = results["df"]
        assignments = results["assignments"]
        model_choice = results["model_choice"]

        tfidf_vectorizer, count_vectorizer, lda_model, nmf_model, _, _ = load_models()
        if model_choice == "NMF":
            model, feature_names = nmf_model, tfidf_vectorizer.get_feature_names_out()
        else:
            model, feature_names = lda_model, count_vectorizer.get_feature_names_out()

        dist = pd.Series(assignments).value_counts().sort_index()

        st.divider()
        st.info(generate_executive_summary(dist, len(result_df), model_choice))

        st.subheader("Topic distribution for this file")
        fig_u, ax_u = plt.subplots(figsize=(8, 3.2))
        themed_axes(fig_u, ax_u)
        bar_colors = [UI_TOPIC_COLORS[i % len(UI_TOPIC_COLORS)] for i in dist.index]
        ax_u.bar([get_topic_label(i) for i in dist.index], dist.values, color=bar_colors)
        ax_u.set_ylabel("Documents")
        plt.setp(ax_u.get_xticklabels(), rotation=25, ha="right")
        ax_u.spines[["top", "right"]].set_visible(False)
        st.pyplot(fig_u)

        st.subheader("Sample assignments")
        st.dataframe(result_df.head(20), width="stretch")

        st.subheader("Download report")
        st.caption("Both formats include the same distribution chart, executive summary, and topic keywords.")

        col_pdf, col_pptx = st.columns(2)
        with col_pdf:
            if st.button("📄 Generate PDF Report", width="stretch"):
                with st.spinner("Building PDF..."):
                    st.session_state["report_bytes"] = build_pdf_report(
                        result_df, assignments, results["topic_dist"],
                        model, feature_names, model_choice, results["source_name"],
                    )
            if "report_bytes" in st.session_state:
                st.download_button(
                    "⬇️ Download PDF",
                    data=st.session_state["report_bytes"],
                    file_name="complaint_topic_report.pdf",
                    mime="application/pdf",
                    width="stretch",
                )

        with col_pptx:
            if st.button("🖼️ Generate Slide Deck", width="stretch"):
                with st.spinner("Building slide deck..."):
                    st.session_state["deck_bytes"] = build_pptx_report(
                        result_df, assignments, results["topic_dist"],
                        model, feature_names, model_choice, results["source_name"],
                    )
            if "deck_bytes" in st.session_state:
                st.download_button(
                    "⬇️ Download Slides (PPTX)",
                    data=st.session_state["deck_bytes"],
                    file_name="complaint_topic_deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    width="stretch",
                )
